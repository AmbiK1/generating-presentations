import os
from telegram import Update
from telegram.ext import Application, CommandHandler, MessageHandler, filters, ContextTypes, ConversationHandler
from pptx import Presentation
from pptx.util import Inches, Pt
import google.generativeai as genai
from google.generativeai import GenerativeModel
import stability_sdk.interfaces.gooseai.generation.generation_pb2 as generation
from stability_sdk import client
from io import BytesIO

# Константы для состояний разговора
TEMA, SLIDES_COUNT, LANGUAGE = range(3)

# API ключи (заменить на свои)
TELEGRAM_TOKEN = 'your-telegram-token'  # Получить у @BotFather в Telegram
GOOGLE_API_KEY = 'your-google-api-key'  # Получить на https://makersuite.google.com/app/apikey
STABILITY_KEY = 'your-stability-key'  # Получить на https://platform.stability.ai/

# Инициализация Gemini
genai.configure(api_key=GOOGLE_API_KEY)
model = GenerativeModel('gemini-pro')

# Инициализация Stability AI
stability_api = client.StabilityInference(
    key=STABILITY_KEY,
    verbose=True,
)


async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text(
        "Привет! Я помогу вам создать презентацию. Пожалуйста, укажите тему презентации:"
    )
    return TEMA


async def get_tema(update: Update, context: ContextTypes.DEFAULT_TYPE):
    context.user_data['tema'] = update.message.text
    await update.message.reply_text("Сколько слайдов должно быть в презентации? (от 3 до 10)")
    return SLIDES_COUNT


async def get_slides_count(update: Update, context: ContextTypes.DEFAULT_TYPE):
    try:
        slides_count = int(update.message.text)
        if 3 <= slides_count <= 10:
            context.user_data['slides_count'] = slides_count
            await update.message.reply_text("На каком языке создать презентацию? (например: русский, английский)")
            return LANGUAGE
        else:
            await update.message.reply_text("Пожалуйста, введите число от 3 до 10")
            return SLIDES_COUNT
    except ValueError:
        await update.message.reply_text("Пожалуйста, введите корректное число")
        return SLIDES_COUNT


async def get_language(update: Update, context: ContextTypes.DEFAULT_TYPE):
    context.user_data['language'] = update.message.text
    await update.message.reply_text("Создаю презентацию, пожалуйста подождите...")

    # Создаем презентацию
    await create_presentation(update, context)
    return ConversationHandler.END


async def create_presentation(update: Update, context: ContextTypes.DEFAULT_TYPE):
    tema = context.user_data['tema']
    slides_count = context.user_data['slides_count']
    language = context.user_data['language']

    # Создаем новую презентацию
    prs = Presentation()

    try:
        for i in range(slides_count):
            await update.message.reply_text(f"Создаю слайд {i + 1}...")

            # Генерируем текст через Gemini
            prompt = f"Create content for slide {i + 1} of {slides_count} about {tema} in {language}. Make it concise and informative."
            response = model.generate_content(prompt)
            slide_content = response.text

            # Создаем изображение через Stability AI
            answers = stability_api.generate(
                prompt=f"Create a professional presentation slide image about {tema}, related to: {slide_content}",
                height=512,
                width=512,
                samples=1,
            )

            # Сохраняем изображение
            image_file = BytesIO()
            for resp in answers:
                for artifact in resp.artifacts:
                    if artifact.type == generation.ARTIFACT_IMAGE:
                        image_file.write(artifact.binary)
            image_file.seek(0)

            # Добавляем слайд
            slide = prs.slides.add_slide(prs.slide_layouts[5])

            # Добавляем текст
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
            tf = txBox.text_frame
            tf.text = slide_content

            # Добавляем изображение
            slide.shapes.add_picture(image_file, Inches(1), Inches(3), Inches(8), Inches(4))

        # Сохраняем презентацию
        presentation_path = f"presentation_{update.effective_user.id}.pptx"
        prs.save(presentation_path)

        # Отправляем файл пользователю
        await update.message.reply_document(document=open(presentation_path, 'rb'))

        # Удаляем временный файл
        os.remove(presentation_path)

    except Exception as e:
        await update.message.reply_text(f"Произошла ошибка при создании презентации: {str(e)}")


def main():
    application = Application.builder().token(TELEGRAM_TOKEN).build()

    conv_handler = ConversationHandler(
        entry_points=[CommandHandler('start', start)],
        states={
            TEMA: [MessageHandler(filters.TEXT & ~filters.COMMAND, get_tema)],
            SLIDES_COUNT: [MessageHandler(filters.TEXT & ~filters.COMMAND, get_slides_count)],
            LANGUAGE: [MessageHandler(filters.TEXT & ~filters.COMMAND, get_language)],
        },
        fallbacks=[],
    )

    application.add_handler(conv_handler)
    application.run_polling()


if __name__ == '__main__':
    main()